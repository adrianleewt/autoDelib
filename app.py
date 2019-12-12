from pptx import Presentation
from pptx.util import Inches
import os
import pandas as pd
from PIL import Image

data = pd.read_csv('Input.csv', encoding = "ISO-8859-1")

prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
ileft = itop = Inches(1)

try:
    dirname = os.path.dirname(os.path.abspath(__file__))
    os.chdir(os.path.join(dirname, 'headshot'))
except:
    print("There is an issue with navigating to the headshot folder. Please make sure the folder exists in the main directory.")

for index,row in data.iterrows():

    slide = prs.slides.add_slide(blank_slide_layout)
    NetID = str(row['NetID'])
    end = NetID.find('@')
    NetID = NetID[:end]

    try:
        img_path = NetID + '.jpg'
        im = Image.open(img_path)
    except:
        try:
            img_path = NetID + '.png'
            im = Image.open(img_path)
        except:
            try:
                img_path = NetID + '.jpeg'
                im = Image.open(NetID + '.jpeg')
            except:
                img_path = NetID + '.JPG'
                im = Image.open(img_path)



    size = im.size

    if size[0] > size[1]:
        fac = size[0]/4
        height = size[1] / fac

        iheight = Inches(height)
        iwidth = Inches(4)

    else:
        fac = size[1]/4
        width = size[0] / fac

        iheight = Inches(4)
        iwidth = Inches(width)

    try:
        pic = slide.shapes.add_picture(img_path,ileft,itop,iwidth,iheight)
    except:
        print(NetID + ' image unsupported. skipping...')


    tleft = Inches(6)
    twidth = Inches(4)
    theight = Inches(3)
    ttop = Inches(4)

    txBox = slide.shapes.add_textbox(tleft,ttop,twidth,theight)
    tf = txBox.text_frame
    tf.word_wrap = True


    tf.text = str(row['Name'])
    # tf.text = row['NetID'] + '\n' + row['Name'] + '\n' + row['Major'] + '\n' + str(row['Year']) + '\nRound 1: ' + str(row['Round 1']) + '\nRound 2: ' + str(row['Round 2'])

    # slide = prs.slides.add_slide(blank_slide_layout)

    # t2left = Inches(1)
    # t2width = Inches(8)
    # t2height = Inches(8)
    # t2top = Inches(1)
    #
    # txBox2 = slide.shapes.add_textbox(t2left,t2top,t2width,t2height)
    # tf2 = txBox2.text_frame
    # tf2.word_wrap = True
    #
    # tf2.text = 'Why GCC: ' + str(row['Why GCC']) + '\n\nRole of New Member: ' + str(row['Role of New Member'])


os.chdir(dirname)
prs.save('test.pptx')
